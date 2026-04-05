from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
OUTPUT_FILE = BASE_DIR / "CS432_Track1_Submission_Presentation_15min.pptx"


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(246, 249, 253)


def add_header(slide, title, subtitle):
    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.5)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(18, 45, 82)
    banner.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.16), Inches(12.2), Inches(0.75))
    ttf = title_box.text_frame
    ttf.clear()
    p = ttf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Calibri"
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(12.2), Inches(0.42))
    stf = subtitle_box.text_frame
    stf.clear()
    p2 = stf.paragraphs[0]
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.name = "Calibri"
    r2.font.size = Pt(19)
    r2.font.color.rgb = RGBColor(224, 232, 244)


def add_notes(slide, notes):
    slide.notes_slide.notes_text_frame.text = notes


def add_bullet_slide(prs, title, subtitle, bullets, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title, subtitle)

    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.9), Inches(11.9), Inches(4.9))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(28, 33, 40)

    add_notes(slide, notes)
    return slide


def add_image_slide(prs, title, subtitle, img_path, caption, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title, subtitle)

    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.9), Inches(1.95), width=Inches(11.5))
    else:
        missing = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.9), Inches(1.95), Inches(11.5), Inches(4.6)
        )
        missing.fill.solid()
        missing.fill.fore_color.rgb = RGBColor(250, 236, 236)
        missing.line.color.rgb = RGBColor(188, 86, 86)
        tf = missing.text_frame
        tf.text = f"Image missing: {img_path.name}"
        tf.paragraphs[0].font.name = "Calibri"
        tf.paragraphs[0].font.size = Pt(22)

    cap = slide.shapes.add_textbox(Inches(0.9), Inches(6.5), Inches(11.6), Inches(0.4))
    ctf = cap.text_frame
    ctf.clear()
    p = ctf.paragraphs[0]
    run = p.add_run()
    run.text = caption
    run.font.name = "Calibri"
    run.font.size = Pt(15)
    run.font.color.rgb = RGBColor(75, 75, 75)

    add_notes(slide, notes)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_bullet_slide(
        prs,
        "CS432 Track 1 Assignment 2",
        "BlindDrop Combined Submission | 15-Minute Viva",
        [
            "Submission folder: CS432_Track1_Submission",
            "Modules: Module A (B+ Tree) + Module B (Secure DB-backed web app)",
            "Presentation objective: demonstrate end-to-end technical completeness",
        ],
        "Time: 0:45\nIntroduce project title, package name, and viva goal. Keep this very short.",
    )

    add_bullet_slide(
        prs,
        "Agenda",
        "What Will Be Covered",
        [
            "Problem model and architecture",
            "Module A implementation, integration, and benchmark proof",
            "Module B authentication, RBAC, security, and SQL optimization",
            "Demo flow and final outcomes",
        ],
        "Time: 0:30\nSet expectations for sequence and mention this is an evidence-backed walkthrough.",
    )

    add_bullet_slide(
        prs,
        "Problem and Domain",
        "BlindDrop Security Model",
        [
            "BlindDrop provides temporary secure file sharing through expiring vaults.",
            "outerToken is public for vault discovery; innerToken is private for authorization.",
            "MAIN token enables full access; SUB token provides restricted file visibility.",
            "Both modules are built on this real domain instead of synthetic sample data.",
        ],
        "Time: 1:30\nExplain why using the real product domain makes the assignment stronger.",
    )

    add_bullet_slide(
        prs,
        "System Architecture",
        "Two Runtime Layers Working Together",
        [
            "Node.js + Express + MySQL + frontend power the product and Module B features.",
            "Python layer provides standalone B+ Tree, integration scripts, and benchmarks.",
            "Google Drive stores file bytes; MySQL stores metadata, sessions, and integrity data.",
            "Evidence package includes docs, logs, benchmark charts, JSON results, and reports.",
        ],
        "Time: 1:30\nShow that architecture is coherent and traceable to submitted files.",
    )

    add_bullet_slide(
        prs,
        "Module A Implementation",
        "From-Scratch B+ Tree",
        [
            "Core features: insert, exact search, update, delete, range query, leaf traversal.",
            "Baseline comparator included: brute-force implementation for fair comparison.",
            "DB/table wrapper demonstrates realistic structured usage of the tree.",
            "Integration indexes real BlindDrop snapshot for project-specific lookup paths.",
        ],
        "Time: 2:00\nEmphasize assignment compliance first, then explain real-domain integration value.",
    )

    add_image_slide(
        prs,
        "Module A Benchmark Evidence",
        "Detailed Speedup Profile",
        BASE_DIR / "Module_A" / "evidence" / "benchmark_speedup.png",
        "Module_A/evidence/benchmark_speedup.png",
        "Time: 1:30\nState headline gains: outer 17.6x, expiry 36.7x, file 62.0x, auth 8.4x.",
    )

    add_bullet_slide(
        prs,
        "Module A Reliability Story",
        "Parity, Rollback, and Rebuild",
        [
            "Parity demo validates index and authoritative state remain synchronized.",
            "Injected failures trigger rollback behavior to prevent silent divergence.",
            "Repair and full rebuild paths prove recovery if mismatch is detected.",
            "Result: performance improvement without sacrificing correctness controls.",
        ],
        "Time: 1:00\nKeep this practical: discuss failure handling, not only happy-path performance.",
    )

    add_bullet_slide(
        prs,
        "Module B Application Layer",
        "Authentication, Session Validation, RBAC",
        [
            "Login and session validation APIs: /api/auth/login and /api/auth/isAuth.",
            "Role mapping follows product model: MAIN -> admin, SUB -> user.",
            "Protected project-specific CRUD on portfolio_entries via /api/portfolio.",
            "Frontend includes member portfolio panel for assignment-facing workflows.",
        ],
        "Time: 2:00\nWalk through one admin flow and one restricted user flow.",
    )

    add_bullet_slide(
        prs,
        "Module B Security Controls",
        "Audit Logging and Tamper Detection",
        [
            "Audit logger records create/update/delete, denied actions, and security checks.",
            "Integrity hash model detects unauthorized direct DB modifications.",
            "Admin-only endpoint /api/security/unauthorized-check surfaces tampered rows.",
            "Normal reads also block tampered entries to enforce integrity at runtime.",
        ],
        "Time: 1:30\nConnect this to evidence folders and explain why this meets assignment security goals.",
    )

    add_image_slide(
        prs,
        "Module B SQL Optimization",
        "Benchmark and EXPLAIN Evidence",
        BASE_DIR / "Module_B" / "evidence" / "BENCHMARK_EVIDENCE" / "03_duration_comparison.png",
        "Module_B/evidence/BENCHMARK_EVIDENCE/03_duration_comparison.png",
        "Time: 1:30\nQuote measured values: 452.8318ms -> 40.0727ms -> 36.8205ms; mention key selected in EXPLAIN.",
    )

    add_bullet_slide(
        prs,
        "Recommended Live Demo Flow",
        "15-Minute Execution Plan",
        [
            "Show health + token model + vault context briefly.",
            "Run/describe Module A index + parity demo and benchmark outcomes.",
            "Demonstrate Module B login, isAuth, and role-based portfolio behavior.",
            "Finish with unauthorized-check and SQL optimization evidence summary.",
        ],
        "Time: 1:30\nUse this as your script if panel asks for live sequence instead of static evidence.",
    )

    add_bullet_slide(
        prs,
        "Final Outcomes",
        "Submission Completeness",
        [
            "Both modules satisfy assignment requirements in one coherent project package.",
            "Artifacts include source, schema, notebook reports, logs, benchmark plots, and docs.",
            "Presentation deck file: CS432_Track1_Submission_Presentation_15min.pptx",
            "Ready for evaluator Q&A and file-path-based verification.",
        ],
        "Time: 0:45\nClose confidently, then move to questions and file proof on demand.",
    )

    prs.save(OUTPUT_FILE)
    print(f"Created: {OUTPUT_FILE}")


if __name__ == "__main__":
    build()

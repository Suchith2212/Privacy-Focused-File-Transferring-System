from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


BASE = Path(__file__).resolve().parent
OUT = BASE / "CS432_Track1_Submission_Presentation_Pro.pptx"


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(246, 249, 253)


def add_header(slide, title, subtitle):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.45)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(16, 42, 79)
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.52), Inches(0.14), Inches(12.2), Inches(0.7))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Calibri"
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)

    sub_box = slide.shapes.add_textbox(Inches(0.52), Inches(0.9), Inches(12.2), Inches(0.34))
    stf = sub_box.text_frame
    stf.clear()
    p2 = stf.paragraphs[0]
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.name = "Calibri"
    r2.font.size = Pt(18)
    r2.font.color.rgb = RGBColor(220, 230, 245)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def bullet_slide(prs, title, subtitle, bullets, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, title, subtitle)

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.68), Inches(1.86), Inches(12.0), Inches(5.28)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = RGBColor(202, 214, 232)

    tf = panel.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(21)
        p.font.color.rgb = RGBColor(25, 30, 38)

    if notes:
        add_notes(slide, notes)
    return slide


def image_slide(prs, title, subtitle, image_path, caption, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, title, subtitle)

    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.78), Inches(1.83), width=Inches(11.95))
    else:
        missing = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.78), Inches(1.83), Inches(11.95), Inches(4.95)
        )
        missing.fill.solid()
        missing.fill.fore_color.rgb = RGBColor(250, 236, 236)
        missing.line.color.rgb = RGBColor(182, 80, 80)
        mtf = missing.text_frame
        mtf.text = f"Image missing: {image_path.name}"
        mtf.paragraphs[0].font.name = "Calibri"
        mtf.paragraphs[0].font.size = Pt(22)

    cap = slide.shapes.add_textbox(Inches(0.78), Inches(6.83), Inches(11.95), Inches(0.36))
    ctf = cap.text_frame
    ctf.clear()
    cp = ctf.paragraphs[0]
    cr = cp.add_run()
    cr.text = caption
    cr.font.name = "Calibri"
    cr.font.size = Pt(14)
    cr.font.color.rgb = RGBColor(84, 84, 84)

    if notes:
        add_notes(slide, notes)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    bullet_slide(
        prs,
        "BlindDrop",
        "CS432 Track 1 Submission | End-to-End Professional Walkthrough",
        [
            "Combined package: Module A (Custom B+ Tree) + Module B (Secure DB-backed Web App)",
            "Submission root: Project_Assignments/CS432_Track1_Submission",
            "Audience: Recruiters, Professors, Interviewers, and Technical Evaluators",
        ],
        "Open with one coherent thesis: real product domain + assignment rigor + evidence-first delivery.",
    )

    bullet_slide(
        prs,
        "Problem Statement",
        "Why This Project Matters",
        [
            "Temporary file sharing needs strict privacy, scoped authorization, and expiry control.",
            "Brute-force metadata access paths do not scale for lookup and range-heavy operations.",
            "Security-critical CRUD requires verifiable RBAC, auditability, and tamper detection.",
            "Optimization claims must be backed by EXPLAIN plans and measured evidence.",
        ],
    )

    bullet_slide(
        prs,
        "Motivation",
        "Build One Defensible Story Instead of Isolated Demos",
        [
            "Reuse real BlindDrop domain rather than disconnected classroom-only examples.",
            "Map both modules to the same token/vault semantics for architectural consistency.",
            "Demonstrate algorithmic value (Module A) and application security value (Module B).",
            "Package source, benchmarks, logs, docs, and runbooks for evaluator verification.",
        ],
    )

    bullet_slide(
        prs,
        "Feature Set",
        "Core + Advanced Capabilities",
        [
            "Vault creation, file upload/list/download, expiry-aware lifecycle, and scoped SUB tokens.",
            "Module A: from-scratch B+ Tree + brute-force baseline + domain-path benchmarking.",
            "Module B: session auth, RBAC portfolio CRUD, audit chain logging, integrity checks.",
            "Adaptive security: CAPTCHA + rate limiting + temporary blocking + evidence endpoint.",
        ],
    )

    bullet_slide(
        prs,
        "Tech Stack",
        "Why These Technologies Were Chosen",
        [
            "Python for custom data structures, benchmark harnesses, and Graphviz rendering.",
            "Node.js + Express for modular API routing and fast backend iteration.",
            "MySQL/InnoDB as authoritative relational store with FK and index-based optimization.",
            "Google Drive for blob storage while MySQL stores access and security metadata.",
            "Vanilla frontend for transparent API-driven demo and easy local evaluator setup.",
        ],
    )

    bullet_slide(
        prs,
        "Architecture",
        "Middleware-Centric Security Data Flow",
        [
            "Frontend -> Express router -> authSession middleware chain -> protected handlers.",
            "requireAuth/requireAdmin revalidate vault and token state on every protected request.",
            "Security is enforced continuously at request time, not only during login issuance.",
            "Route handlers coordinate MySQL authority, audit logs, and Drive blob operations.",
            "Module A uses exported DB snapshots for indexing/parity/benchmark workflows.",
        ],
        "Call out the differentiator: session token is rechecked against DB state for each protected route.",
    )

    bullet_slide(
        prs,
        "Database Schema",
        "Relational Backbone of the System",
        [
            "Core tables: vaults, inner_tokens, files, file_metadata, file_key_access, sessions, auth_attempts, download_logs, captcha_tracking, expiry_jobs.",
            "Assignment table: portfolio_entries with owner scope + integrity hash for tamper checks.",
            "Indexed paths include token lookup, RBAC listing, integrity, and audit timeline access.",
            "Schema is rerunnable and submission-portable through Module_B/sql/init_schema.sql.",
        ],
        "Explain that docs sometimes refer to 9-table core; packaged schema includes extended operational tables plus portfolio.",
    )

    bullet_slide(
        prs,
        "Folder Structure",
        "What Evaluators Should Inspect First",
        [
            "Docs/: architecture notes, API reference, demo guide, and submission checklist.",
            "Module_A/: database core, integration layer, render outputs, benchmark evidence.",
            "Module_B/app/backend/: routes, services, SQL, verification and benchmark reports.",
            "Module_B/app/frontend/: full UI flow for vault/file/member/portfolio operations.",
            "Module_B/evidence/: API, DB, audit, and benchmark proof artifacts.",
        ],
    )

    bullet_slide(
        prs,
        "Deep Dive",
        "Core Modules and Design Decisions",
        [
            "BPlusTree supports split/merge/borrow balancing and linked-leaf range traversal.",
            "PostingListIndex enables duplicate-key one-to-many mapping without rewriting tree core.",
            "Parity manager proves rollback, divergence detection, lazy repair, and rebuild behavior.",
            "Upload security pipeline includes fileValidation.js and filePathSchema.js safeguards.",
            "portfolioIntegrity.js + security.js enforce tamper checks and adaptive abuse controls.",
        ],
    )

    bullet_slide(
        prs,
        "Execution Flow",
        "Input -> Processing -> Output",
        [
            "Create vault and upload files -> generate outer token + MAIN inner token.",
            "Access flow verifies outer + inner token and resolves MAIN/SUB effective permissions.",
            "Portfolio APIs apply session auth + RBAC + integrity checks before CRUD actions.",
            "Download flow validates access, logs event, and enforces one-time file consumption.",
            "Module A pipeline indexes exported snapshot and emits benchmark/visual evidence.",
        ],
    )

    bullet_slide(
        prs,
        "Challenges and Resolutions",
        "What Was Hard and How It Was Solved",
        [
            "Duplicate logical keys in domain indexes -> solved with posting-list abstraction.",
            "Potential index-authority drift -> solved with parity validation + repair + rebuild.",
            "Hashed credential lookup performance -> solved with token_lookup_hash prefilter.",
            "Portable tamper defense across MySQL environments -> app-side integrity strategy.",
            "Abuse resistance -> weighted failures, CAPTCHA lifecycle, and adaptive blocking.",
        ],
    )

    bullet_slide(
        prs,
        "Performance Highlights",
        "Measured, Not Claimed",
        [
            "Module A domain benchmark: strong speedups across lookup and range paths.",
            "Module B SQL benchmark: 452.8318 ms full scan -> 40.0727 ms indexed lookup.",
            "Memory at 21,000 records: B+ Tree retains ~1,239 KB vs brute-force ~1,892 KB.",
            "Module A evidence includes 20-point domain + 22-point detailed benchmark suites.",
        ],
    )

    image_slide(
        prs,
        "Module A Evidence",
        "Speedup Benchmark Dashboard",
        BASE / "Module_A" / "evidence" / "benchmark_speedup.png",
        "Source: Module_A/evidence/benchmark_speedup.png",
    )

    image_slide(
        prs,
        "Module B Evidence",
        "SQL Duration Comparison",
        BASE / "Module_B" / "evidence" / "BENCHMARK_EVIDENCE" / "03_duration_comparison.png",
        "Source: Module_B/evidence/BENCHMARK_EVIDENCE/03_duration_comparison.png",
    )

    bullet_slide(
        prs,
        "Demo Walkthrough",
        "Evaluator-Friendly Verification Sequence",
        [
            "Live flow: health -> admin login -> isAuth -> portfolio CRUD -> denied user action -> unauthorized-check.",
            "Module A scripts: blinddrop_index_demo.py, db_index_parity_demo.py, benchmark scripts.",
            "Module B scripts: seed_module_b_demo.js, verify_module_b_e2e.js, index_benchmark.js.",
            "Video references: youtu.be/T24vXjLI5dI (Module A), youtu.be/FzY8OeX4d5E (Module B).",
        ],
    )

    bullet_slide(
        prs,
        "Limitations and Roadmap",
        "Current Constraints and Next Steps",
        [
            "Session/security state is in-memory; restart resets adaptive controls.",
            "Helper SUB-token secret storage is practical for demo but should be hardened.",
            "Future: Redis-backed session/security state, immutable external audit sink, CI pipeline.",
            "Future: deeper upload content scanning and stronger key lifecycle governance.",
        ],
    )

    bullet_slide(
        prs,
        "Conclusion",
        "Impact and Final Positioning",
        [
            "One coherent submission unifies algorithmic rigor and secure application engineering.",
            "Module A proves index performance/correctness under realistic domain paths.",
            "Module B proves enforceable RBAC, tamper detection, and measurable optimization.",
            "Evidence-first packaging makes this project audit-friendly and interview-ready.",
        ],
    )

    bullet_slide(
        prs,
        "Q&A",
        "Thank You",
        [
            "Questions on architecture, security tradeoffs, indexing, or benchmark methodology.",
            "Submission path: Project_Assignments/CS432_Track1_Submission",
            "Deck file: CS432_Track1_Submission_Presentation_Pro.pptx",
        ],
    )

    prs.save(OUT)
    print(f"Created: {OUT}")


if __name__ == "__main__":
    build()

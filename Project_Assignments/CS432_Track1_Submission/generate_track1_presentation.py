from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
OUTPUT_FILE = BASE_DIR / "CS432_Track1_Submission_Presentation.pptx"


def style_title(run):
    run.font.name = "Calibri"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)


def style_subtitle(run):
    run.font.name = "Calibri"
    run.font.size = Pt(22)
    run.font.color.rgb = RGBColor(230, 230, 230)


def add_banner(slide, title_text, subtitle_text):
    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.6)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(22, 43, 72)
    banner.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.clear()
    p = title_tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    style_title(run)

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.45))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.clear()
    p = subtitle_tf.paragraphs[0]
    run = p.add_run()
    run.text = subtitle_text
    style_subtitle(run)


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 248, 252)


def add_bullet_slide(prs, title, subtitle, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_banner(slide, title, subtitle)

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.95), Inches(11.9), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(25)
        p.font.color.rgb = RGBColor(30, 30, 30)

    return slide


def add_metric_slide(prs, title, subtitle, metric_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_banner(slide, title, subtitle)

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.7)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = RGBColor(180, 195, 214)

    tf = panel.text_frame
    tf.clear()
    tf.word_wrap = True

    for idx, line in enumerate(metric_lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(26 if idx == 0 else 23)
        p.font.bold = idx == 0
        p.font.color.rgb = RGBColor(20, 33, 54)

    return slide


def add_image_slide(prs, title, subtitle, image_path, caption):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_banner(slide, title, subtitle)

    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.9), Inches(2.0), width=Inches(11.5))
    else:
        missing = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.3)
        )
        missing.fill.solid()
        missing.fill.fore_color.rgb = RGBColor(250, 236, 236)
        missing.line.color.rgb = RGBColor(199, 96, 96)
        t = missing.text_frame
        t.text = f"Image missing: {image_path.name}"
        t.paragraphs[0].font.name = "Calibri"
        t.paragraphs[0].font.size = Pt(22)
        t.paragraphs[0].font.color.rgb = RGBColor(130, 40, 40)

    cap_box = slide.shapes.add_textbox(Inches(0.9), Inches(6.45), Inches(11.5), Inches(0.55))
    cap_tf = cap_box.text_frame
    cap_tf.clear()
    p = cap_tf.paragraphs[0]
    run = p.add_run()
    run.text = caption
    run.font.name = "Calibri"
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(60, 60, 60)

    return slide


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_bullet_slide(
        prs,
        "CS432 Track 1 Submission",
        "BlindDrop End-to-End Presentation",
        [
            "Student Project: Privacy-Focused File Transferring System",
            "Submission Package: CS432_Track1_Submission",
            "Modules Covered: Module A (B+ Tree) + Module B (Secure Web App)",
        ],
    )

    add_bullet_slide(
        prs,
        "Agenda",
        "Complete Walkthrough",
        [
            "1) Problem and product model",
            "2) Submission structure and architecture",
            "3) Module A: implementation, parity, benchmarks",
            "4) Module B: auth, RBAC, CRUD, audit, tamper detection",
            "5) SQL optimization evidence and final outcomes",
        ],
    )

    add_bullet_slide(
        prs,
        "Problem Statement",
        "Why BlindDrop Matters",
        [
            "BlindDrop is a secure temporary file sharing system with expiring vaults.",
            "Outer token is public for discovery; inner token is private for authorization.",
            "MAIN token grants full vault access; SUB token supports restricted file scope.",
            "Goal: deliver both assignment modules on top of a real working product domain.",
        ],
    )

    add_bullet_slide(
        prs,
        "Submission Strategy",
        "Single Coherent Package",
        [
            "Module A: custom B+ Tree + brute-force baseline + DB wrapper + integration layer.",
            "Module B: local DB-backed app with session auth, RBAC, CRUD, audit, and SQL profiling.",
            "Evidence-first approach: reports, console logs, JSON snapshots, benchmark charts.",
            "Demo-ready package: docs, runbook, API reference, and hosted demo videos.",
        ],
    )

    add_bullet_slide(
        prs,
        "Architecture Overview",
        "Runtime and Data Layers",
        [
            "Node.js + Express + MySQL + static frontend powers product + Module B workflows.",
            "Python module provides standalone B+ Tree engine and benchmarking/visualization tools.",
            "Google Drive stores file bytes in product flow; MySQL stores metadata and security state.",
            "Module A indexes are built from exported BlindDrop snapshot for real-domain benchmarking.",
        ],
    )

    add_bullet_slide(
        prs,
        "Core Product Flow",
        "BlindDrop End-to-End",
        [
            "Create vault and upload files -> receive outerToken + MAIN innerToken.",
            "Use outerToken + MAIN to access full vault and provision SUB tokens.",
            "Use outerToken + SUB to access only permitted files.",
            "Expiry-aware behavior enforces temporary sharing lifecycle.",
        ],
    )

    add_bullet_slide(
        prs,
        "Module A",
        "From-Scratch B+ Tree + Integration",
        [
            "Standalone engine: insert, search, update, delete, range query, linked leaves.",
            "Baseline comparator: brute-force implementation for performance contrast.",
            "Integration indexes built for outer-token, expiry range, files, and auth timelines.",
            "Parity/rebuild story proves rollback, validation, repair, and authoritative rebuild.",
        ],
    )

    add_metric_slide(
        prs,
        "Module A Benchmark Outcomes",
        "Packaged Evidence Highlights",
        [
            "Domain benchmark average speedups",
            "Outer-token lookup: 17.6x faster",
            "Expiry range scan: 36.7x faster",
            "Vault-file range scan: 62.0x faster",
            "Auth timeline range scan: 8.4x faster",
            "Detailed suite: 22 benchmark points | Domain suite: 20 points",
        ],
    )

    add_image_slide(
        prs,
        "Module A Evidence",
        "Speedup Dashboard",
        BASE_DIR / "Module_A" / "evidence" / "benchmark_speedup.png",
        "Source: Module_A/evidence/benchmark_speedup.png",
    )

    add_image_slide(
        prs,
        "Module A Visualization",
        "Integrated B+ Tree Render",
        BASE_DIR / "Module_A" / "integration" / "bptree_v2" / "01_vaults__outer_token.png",
        "Source: Module_A/integration/bptree_v2/01_vaults__outer_token.png",
    )

    add_bullet_slide(
        prs,
        "Module B",
        "Application and Security Layer",
        [
            "Local Express + MySQL app with session-backed authentication.",
            "Role mapping reuses product credentials: MAIN -> admin, SUB -> user.",
            "Project-specific CRUD surface: portfolio_entries via /api/portfolio.",
            "Frontend includes Member Portfolio panel for assignment workflow visibility.",
        ],
    )

    add_bullet_slide(
        prs,
        "Module B Security",
        "RBAC, Audit, Tamper Detection",
        [
            "Session middleware protects portfolio and security routes.",
            "Audit events recorded for create/update/delete/denied/security actions.",
            "Integrity hash model detects unauthorized direct DB modifications.",
            "Admin-only endpoint /api/security/unauthorized-check reports tampered rows.",
        ],
    )

    add_metric_slide(
        prs,
        "Module B SQL Optimization",
        "Measured Benchmark Results",
        [
            "Protected portfolio query benchmark",
            "Baseline full scan: 452.8318 ms",
            "Composite lookup index: 40.0727 ms (11.30x faster)",
            "Composite + covering comparison: 36.8205 ms (12.30x faster)",
            "Captured EXPLAIN still uses idx_portfolio_benchmark_lookup",
        ],
    )

    add_image_slide(
        prs,
        "Module B Benchmark Evidence",
        "Duration Comparison",
        BASE_DIR / "Module_B" / "evidence" / "BENCHMARK_EVIDENCE" / "03_duration_comparison.png",
        "Source: Module_B/evidence/BENCHMARK_EVIDENCE/03_duration_comparison.png",
    )

    add_bullet_slide(
        prs,
        "Live Demo Sequence",
        "Recommended Presentation Order",
        [
            "1) Explain vault/token model and show application health.",
            "2) Demonstrate MAIN and SUB token access behavior.",
            "3) Run Module A index + parity demos and show benchmark dashboards.",
            "4) Show Module B login/isAuth/portfolio CRUD and RBAC restrictions.",
            "5) Trigger unauthorized-check and present SQL benchmark + EXPLAIN evidence.",
        ],
    )

    add_bullet_slide(
        prs,
        "Deliverables Recap",
        "What Was Submitted",
        [
            "Complete source code, schema, reports, notebooks, and evidence folders.",
            "19 B+ Tree render images + benchmark dashboards + JSON summaries.",
            "Audit log and tamper detection proof integrated with application APIs.",
            "Demo videos: Module A and Module B links included in README/report files.",
        ],
    )

    add_bullet_slide(
        prs,
        "Thank You",
        "Q&A",
        [
            "Submission folder: Project_Assignments/CS432_Track1_Submission",
            "Deck file: CS432_Track1_Submission_Presentation.pptx",
            "Prepared for end-to-end evaluator walkthrough.",
        ],
    )

    prs.save(OUTPUT_FILE)
    print(f"Presentation created: {OUTPUT_FILE}")


if __name__ == "__main__":
    build_presentation()
